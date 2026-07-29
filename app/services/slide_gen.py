"""PPT (.pptx) generation service with AI-powered backgrounds.

Generates a real PowerPoint file from a script, with AI-generated
background images, professional layouts, and elegant typography.

Free image models supported:
  - Pollinations.ai (Flux, no API key needed)
  - Google Gemini 2.5 Flash Image (500 free images/day)
"""

import asyncio
import base64
import json
import logging
import os
import textwrap
from pathlib import Path
from urllib.parse import quote

import httpx
from fastapi import HTTPException
from openai import AsyncOpenAI
from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from app.config import settings

logger = logging.getLogger(__name__)

STORAGE_ROOT = Path("storage")
SLIDES_DIR = STORAGE_ROOT / "slides"

# PPT dimensions (16:9 widescreen)
PPT_WIDTH_EMU = 12192000   # ~13.33 inches
PPT_HEIGHT_EMU = 6858000   # ~7.5 inches

# Inches helpers
IN = Inches

# ===========================================================================
# Rich Image Style System
# Each style has: prompt_modifier (injected into AI prompts),
# negative (what to avoid), visual_label (for UI), visual_icon (emoji),
# description (tooltip for UI), palette_hint (color direction)
# ===========================================================================

STYLE_METADATA = {
    "realistic": {
        "label": "写实摄影",
        "icon": "📷",
        "desc": "电影级实拍质感，光影细腻，适合需要真实感的场景",
    },
    "illustration": {
        "label": "现代插画",
        "icon": "🎨",
        "desc": "扁平化数字插画，色彩明快，适合知识科普类内容",
    },
    "anime": {
        "label": "动漫风",
        "icon": "🌸",
        "desc": "日系动漫风格，明快鲜艳，适合故事/情感类内容",
    },
    "watercolor": {
        "label": "水彩手绘",
        "icon": "🖌️",
        "desc": "柔和水彩质感，清新文艺，适合散文/随笔类内容",
    },
    "comic": {
        "label": "漫画分镜",
        "icon": "💥",
        "desc": "美漫风格，粗线条强对比，适合观点鲜明的内容",
    },
    "minimal": {
        "label": "极简商务",
        "icon": "◻️",
        "desc": "简约几何构图，高级灰调，适合商业/科技类PPT",
    },
    "cinematic": {
        "label": "电影质感",
        "icon": "🎬",
        "desc": "宽银幕电影色调，景深感强，适合叙事类内容",
    },
    "dreamy": {
        "label": "梦幻柔焦",
        "icon": "✨",
        "desc": "柔和光晕效果，温暖的梦幻氛围，适合情感/治愈类",
    },
    "vintage": {
        "label": "复古胶片",
        "icon": "📽️",
        "desc": "胶片颗粒感，怀旧暖色调，适合历史/回忆类内容",
    },
    "cyberpunk": {
        "label": "赛博朋克",
        "icon": "🌃",
        "desc": "霓虹灯光、科技感未来都市，适合科技/趋势类内容",
    },
    "nature": {
        "label": "自然风光",
        "icon": "🌿",
        "desc": "自然景观、植物纹理、温暖阳光，适合治愈/环保类",
    },
    "abstract": {
        "label": "抽象艺术",
        "icon": "🌀",
        "desc": "流体色彩、几何抽象，适合哲学/创意类内容",
    },
}

# Prompt modifiers injected into AI image generation
SLIDE_STYLE_MAP = {
    "realistic": "photorealistic, cinematic, dramatic natural lighting, ultra detailed textures, 8k resolution, professional photography, depth of field, bokeh background, shot on 85mm lens",
    "illustration": "digital illustration, flat vector art style, clean modern design, vibrant harmonious color palette, professional editorial design, smooth gradients, crisp edges",
    "anime": "anime art style, cel shading, vibrant saturated colors, Japanese animation aesthetic, detailed background art, soft rim lighting, Studio Ghibli inspired composition",
    "watercolor": "watercolor painting on textured paper, soft translucent brush strokes, artistic bleeding edges, pastel muted tones, elegant dreamy atmosphere, impressionist style",
    "comic": "comic book art style, bold ink outlines, dynamic action composition, vivid pop colors, halftone dot shading, graphic novel aesthetic, strong contrast",
    "minimal": "minimalist design, clean geometric composition, smooth subtle gradients, modern corporate aesthetic, negative space, soft diffused lighting, matte finish",
    "cinematic": "cinematic film still, anamorphic lens, golden hour lighting, shallow depth of field, film grain, color graded teal and orange tones, widescreen composition",
    "dreamy": "ethereal dreamy atmosphere, soft glow, lens flare, pastel color palette, hazy morning light, romantic aesthetic, floating particles, cloud-like textures",
    "vintage": "vintage film photography, warm sepia tones, light leaks, film grain texture, nostalgic atmosphere, retro aesthetic, soft faded colors, Kodak Portra 400 look",
    "cyberpunk": "cyberpunk aesthetic, neon lights, rain-slicked streets, holographic reflections, futuristic cityscape, purple and cyan color scheme, Blade Runner atmosphere",
    "nature": "lush nature landscape, golden sunlight through leaves, organic textures, botanical details, macro nature photography, morning dew, biophilic design",
    "abstract": "abstract fluid art, alcohol ink technique, swirling organic shapes, rich jewel tones, marble texture, modern gallery wall aesthetic, sophisticated composition",
}

# Universal negative prompt — applied to all AI image generations
NEGATIVE_PROMPT = (
    "no text, no letters, no numbers, no watermarks, no signatures, "
    "no logos, no people faces in foreground, no clutter, no low quality, "
    "no blurry, no distorted, no ugly, no deformed"
)

# Quality-enhancing suffix appended to every prompt
QUALITY_SUFFIX = "masterpiece, best quality, highly detailed, professional, sharp focus"

# Fallback gradient palettes (when no AI image)
GRADIENT_PALETTES = [
    {"top": (15, 32, 75), "bottom": (30, 80, 160)},
    {"top": (40, 20, 60), "bottom": (100, 50, 120)},
    {"top": (20, 50, 40), "bottom": (30, 100, 80)},
    {"top": (50, 25, 20), "bottom": (140, 70, 40)},
    {"top": (40, 20, 25), "bottom": (120, 40, 60)},
    {"top": (20, 35, 55), "bottom": (40, 90, 130)},
    {"top": (25, 25, 35), "bottom": (60, 50, 90)},
    {"top": (30, 35, 20), "bottom": (70, 90, 40)},
]

# Font path for CJK
CJK_FONT_CANDIDATES = [
    "/System/Library/Fonts/PingFang.ttc",
    "/System/Library/Fonts/STHeiti Light.ttc",
    "/System/Library/Fonts/Hiragino Sans GB.ttc",
    "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
    "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
]

PPTX_FONT_NAME = "Microsoft YaHei"  # PPT standard CJK font


def _ensure_dir() -> None:
    SLIDES_DIR.mkdir(parents=True, exist_ok=True)


def _load_font(size: int):
    """Load a CJK font for Pillow rendering (used for bg preprocessing)."""
    from PIL import ImageFont
    for path in CJK_FONT_CANDIDATES:
        if os.path.exists(path):
            try:
                return ImageFont.truetype(path, size)
            except Exception:
                continue
    return ImageFont.load_default()


# ---------------------------------------------------------------------------
# AI image generation for slide backgrounds
# ---------------------------------------------------------------------------

async def _generate_slide_bg_image(
    prompt: str,
    seg_id: int,
    note_id: int,
    model_config: dict | None = None,
) -> str | None:
    """Generate a background image via AI and return the local file path.

    Priority: free models (Pollinations / Gemini) -> SiliconFlow paid -> None.
    """
    cfg = model_config or {}
    model_id = cfg.get("model_id") or settings.IMAGE_MODEL

    # --- Try free models first (with retry) ---
    if _is_free_model(model_id):
        result = await _retry_image_generation(
            _generate_image_with_free_model_wrapper,
            prompt=prompt, seg_id=seg_id, note_id=note_id,
            max_retries=2, model_config=model_config,
            width=1920, height=1080, prefix="bg",
        )
        if result:
            return result
        logger.warning("Free model %s failed for slide segment %s, using gradient fallback", model_id, seg_id)
        return None

    # --- Existing SiliconFlow / paid model path ---
    api_key = cfg.get("api_key") or settings.IMAGE_API_KEY
    base_url = cfg.get("base_url") or settings.IMAGE_BASE_URL

    if not api_key:
        logger.debug("No IMAGE_API_KEY configured, skipping AI background generation")
        return None

    # Determine optimal inference steps per model
    m = model_id.lower()
    if "schnell" in m:
        steps = 4
    elif "lightning" in m:
        steps = 4
    elif "flux" in m and "dev" in m:
        steps = 28
    elif "flux" in m:
        steps = 20
    elif "kolors" in m:
        steps = 25
    elif "qwen-image" in m or "z-image" in m or "ernie" in m:
        steps = 20
    elif "sd3" in m or "stable-diffusion-3" in m:
        steps = 28
    elif "turbo" in m:
        steps = 8
    elif "playground" in m:
        steps = 25
    else:
        steps = 20

    try:
        url = f"{base_url.rstrip('/')}/images/generations"
        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
        }
        payload = {
            "model": model_id,
            "prompt": prompt,
            "image_size": "1920x1080",
            "num_inference_steps": steps,
        }

        async with httpx.AsyncClient(timeout=120.0) as client:
            resp = await client.post(url, json=payload, headers=headers)
            if resp.status_code in (401, 402, 403):
                try:
                    _msg = resp.json().get("message", "")
                except Exception:
                    _msg = ""
                if resp.status_code == 402 or "insufficient" in _msg.lower() or "余额" in _msg:
                    logger.error("Image model %s returned %s: %s", model_id, resp.status_code, _msg)
                    raise HTTPException(status_code=402, detail=f"图片模型 {model_id} 余额不足，请充值后重试")
                logger.error("Image model %s returned %s: %s", model_id, resp.status_code, _msg)
                raise HTTPException(status_code=resp.status_code, detail=f"图片模型 {model_id} 不可用：{_msg or '模型未开启'}")
            resp.raise_for_status()
            data = resp.json()

        image_url = None
        images = data.get("images", [])
        if images and isinstance(images, list):
            image_url = images[0].get("url")
        elif data.get("data"):
            image_url = data["data"][0].get("url")

        if not image_url:
            logger.warning("No image URL returned for slide bg segment %s", seg_id)
            return None

        out_path = SLIDES_DIR / f"note_{note_id}_bg_{seg_id}.png"
        async with httpx.AsyncClient(timeout=60.0) as client:
            img_resp = await client.get(image_url)
            img_resp.raise_for_status()
            out_path.write_bytes(img_resp.content)

        return str(out_path)

    except HTTPException:
        raise  # Let 403 errors bubble up to caller
    except Exception:
        logger.exception("Failed to generate AI background for slide segment %s", seg_id)
        return None


# ===========================================================================
# Free image generation APIs (Pollinations.ai + Google Gemini)
# ===========================================================================

# Pollinations.ai — completely free, no API key needed for anonymous use
POLLINATIONS_BASE = "https://image.pollinations.ai/prompt"

# Available models on Pollinations.ai
POLLINATIONS_MODELS = {
    "pollinations-flux": "flux",           # Flux (best quality, free)
    "pollinations-turbo": "turbo",        # Turbo (fastest)
    "pollinations-flux-realism": "flux-realism",  # Flux-realism variant
}


async def _generate_image_pollinations(
    prompt: str,
    width: int = 1080,
    height: int = 1440,
    seg_id: int = 0,
    note_id: int = 0,
    model_name: str = "flux",
    prefix: str = "xhs_bg",
) -> str | None:
    """Generate an image using Pollinations.ai free API.

    No API key required for anonymous use. Registering a free account
    (enter.pollinations.ai) raises rate limits and removes watermark.

    Args:
        prompt: Text prompt (English recommended, Chinese also supported)
        width: Output image width
        height: Output image height
        seg_id: Segment ID for file naming
        note_id: Note ID for file naming
        model_name: Model to use (flux, turbo, flux-realism)
        prefix: File prefix for output path

    Returns:
        Local file path to the generated image, or None on failure.
    """
    try:
        # Build enhanced prompt — only inject quality suffix into the text prompt
        enhanced_prompt = f"{prompt}, {QUALITY_SUFFIX}"
        encoded_prompt = quote(enhanced_prompt)

        params = {
            "width": str(width),
            "height": str(height),
            "model": model_name,
            "nologo": "true",
            "enhance": "true",
        }

        # For Flux models, use the dedicated "negative" query parameter
        # instead of polluting the main prompt text (which confuses the model)
        if _has_negative_support(model_name):
            params["negative"] = quote(NEGATIVE_PROMPT)

        # Check if user has configured a Pollinations API key for higher limits
        api_key = getattr(settings, "POLLINATIONS_API_KEY", "") or ""
        if api_key:
            params["key"] = api_key

        param_str = "&".join(f"{k}={v}" for k, v in params.items())
        url = f"{POLLINATIONS_BASE}/{encoded_prompt}?{param_str}"

        async with httpx.AsyncClient(timeout=90.0) as client:
            resp = await client.get(url)
            if resp.status_code != 200:
                logger.warning(
                    "Pollinations.ai returned %d for segment %s: %s",
                    resp.status_code, seg_id, resp.text[:200],
                )
                return None

            out_path = SLIDES_DIR / f"note_{note_id}_{prefix}_{seg_id}.png"
            out_path.write_bytes(resp.content)

        logger.info("Pollinations.ai generated image for segment %s (%dx%d)", seg_id, width, height)
        return str(out_path)

    except Exception:
        logger.exception("Pollinations.ai generation failed for segment %s", seg_id)
        return None


async def _generate_image_gemini(
    prompt: str,
    width: int = 1080,
    height: int = 1440,
    seg_id: int = 0,
    note_id: int = 0,
    prefix: str = "xhs_bg",
    model_config: dict | None = None,
) -> str | None:
    """Generate an image using Google Gemini 2.5 Flash Image API (free tier).

    Free tier: 500 images/day at 1024x1024 max.
    Requires GEMINI_API_KEY in .env or model_config["api_key"]
    (get from aistudio.google.com).

    Uses the official Gemini generateContent API:
    POST https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash-image:generateContent

    Args:
        prompt: Image description (supports Chinese)
        width: Output width (max 1024 on free tier)
        height: Output height (max 1024 on free tier)
        seg_id: Segment ID for file naming
        note_id: Note ID for file naming
        prefix: File prefix for output path
        model_config: Dict with optional "api_key" for Gemini key

    Returns:
        Local file path to the generated image, or None on failure.
    """
    # Try model_config api_key first, then settings.GEMINI_API_KEY
    cfg = model_config or {}
    api_key = cfg.get("api_key") or getattr(settings, "GEMINI_API_KEY", "") or ""
    if not api_key:
        logger.warning("No Gemini API key configured for segment %s", seg_id)
        return None

    # Free tier max is 1024x1024
    max_dim = 1024
    if width > max_dim or height > max_dim:
        scale = max_dim / max(width, height)
        width = int(width * scale)
        height = int(height * scale)

    try:
        url = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash-image:generateContent"
        headers = {
            "Content-Type": "application/json",
            "x-goog-api-key": api_key,
        }

        # Enhance prompt for better results
        style_hint = ""
        if model_config and model_config.get("image_style"):
            st = model_config.get("image_style", "")
            if st in SLIDE_STYLE_MAP:
                style_hint = f", {SLIDE_STYLE_MAP[st][:120]}"

        if any("\u4e00" <= ch <= "\u9fff" for ch in prompt[:20]):
            full_prompt = (
                f"Professional high-quality photograph, {prompt}{style_hint}, "
                f"soft natural lighting, warm aesthetic tones, "
                f"3:4 portrait orientation, {QUALITY_SUFFIX}, {NEGATIVE_PROMPT}"
            )
        else:
            full_prompt = f"{prompt}{style_hint}, {QUALITY_SUFFIX}, {NEGATIVE_PROMPT}"

        body = {
            "contents": [{"parts": [{"text": full_prompt}]}],
            "generationConfig": {"responseModalities": ["IMAGE", "TEXT"]},
        }

        async with httpx.AsyncClient(timeout=120.0) as client:
            resp = await client.post(url, json=body, headers=headers)

            if resp.status_code in (429, 403):
                logger.warning(
                    "Gemini API rate limited or denied for segment %s: %d",
                    seg_id, resp.status_code,
                )
                return None
            if resp.status_code != 200:
                logger.warning(
                    "Gemini API returned %d for seg %s: %s",
                    resp.status_code, seg_id, resp.text[:400],
                )
                return None

            data = resp.json()

        # Extract image from interactions response (handles multiple formats)
        image_data = _extract_gemini_image_data(data)

        if not image_data:
            logger.warning("Gemini returned no image data for segment %s; response keys: %s",
                           seg_id, list(data.keys())[:5])
            return None

        out_path = SLIDES_DIR / f"note_{note_id}_{prefix}_{seg_id}.png"
        out_path.write_bytes(image_data)

        logger.info(
            "Gemini Flash Image generated for segment %s (%dx%d)",
            seg_id, width, height,
        )
        return str(out_path)

    except Exception:
        logger.exception("Gemini generation failed for segment %s", seg_id)
        return None


# Free models registry — each entry defines model_id, provider, and width/height hints
FREE_IMAGE_MODELS = {
    "pollinations-flux": {
        "provider": "pollinations",
        "model_name": "flux",
        "max_width": 1920,
        "max_height": 1440,
        "label": "Flux 写实风格",
        "description": "最推荐的小红书图文模型。画质细腻、光影自然，适合生活方式/穿搭/美食/旅行等主流内容。随手出片，氛围感拉满。",
        "features": ["画质细腻", "光影自然", "氛围感强", "支持负向提示词"],
        "speed": "中",
        "quality": "高",
        "best_for": "生活穿搭 / 美食旅行 / 日常分享",
        "requires_key": False,
        "icon": "🖼️",
    },
    "pollinations-turbo": {
        "provider": "pollinations",
        "model_name": "turbo",
        "max_width": 1920,
        "max_height": 1440,
        "label": "Turbo 快速出图",
        "description": "生成速度最快的模型，适合批量生成和预览迭代。画质稍逊于 Flux，但胜在效率——3-5 秒出图，灵感不等待。",
        "features": ["极速生成", "适合批量", "预览迭代"],
        "speed": "快",
        "quality": "中",
        "best_for": "快速预览 / 批量生成 / 头脑风暴",
        "requires_key": False,
        "icon": "⚡",
    },
    "gemini-flash-image": {
        "provider": "gemini",
        "model_name": "gemini-2.5-flash-image",
        "max_width": 1024,
        "max_height": 1024,
        "label": "Gemini 智能构图",
        "description": "Google 最新的免费图片模型，擅长理解复杂语义和中文提示词。每天 500 张免费额度，需要申请 API Key（无需信用卡）。对构图和色彩的把控很出色。",
        "features": ["理解中文", "语义精准", "构图出色", "每天500张"],
        "speed": "中",
        "quality": "高",
        "best_for": "复杂场景 / 中文语义 / 创意构图",
        "requires_key": True,
        "icon": "🧠",
    },
}


def _is_free_model(model_id: str) -> bool:
    """Check if a model_id refers to a free model."""
    return model_id in FREE_IMAGE_MODELS


def _has_negative_support(model_name: str) -> bool:
    """Check if the model supports negative prompts (Flux does, Turbo may not)."""
    return model_name in ("flux", "flux-realism")


def _extract_gemini_image_data(response_data: dict) -> bytes | None:
    """Extract image bytes from Gemini interactions API response.

    Handles multiple possible response formats robustly:
    1. steps[].content[] with type=image and data field
    2. steps[].content[] with inline_data / inlineData
    3. steps[].output.candidates[] (generateContent embedded)
    4. Top-level candidates[] (old generateContent format)
    5. Base64 padding fixes for stripped trailing '='
    """
    def _try_decode(b64_str: str) -> bytes | None:
        """Try to decode base64, auto-fixing missing padding."""
        try:
            return base64.b64decode(b64_str)
        except Exception:
            missing = len(b64_str) % 4
            if missing:
                b64_str += "=" * (4 - missing)
            try:
                return base64.b64decode(b64_str)
            except Exception:
                return None

    def _extract_from_parts(parts: list) -> bytes | None:
        """Extract image from parts array (generateContent format)."""
        for part in parts:
            for inline_key in ("inlineData", "inline_data"):
                inline = part.get(inline_key)
                if isinstance(inline, dict) and inline.get("data"):
                    return _try_decode(str(inline["data"]))
        return None

    def _extract_from_blocks(blocks: list) -> bytes | None:
        """Extract image from content blocks (interactions format)."""
        for block in blocks:
            if block.get("type") != "image":
                continue
            if block.get("data"):
                return _try_decode(str(block["data"]))
            for inline_key in ("inlineData", "inline_data"):
                inline = block.get(inline_key)
                if isinstance(inline, dict) and inline.get("data"):
                    return _try_decode(str(inline["data"]))
        return None

    # Format 1 & 2: steps[].content[] or steps[].output
    for step in response_data.get("steps", []):
        content = step.get("content") or step.get("summary") or []
        if isinstance(content, list):
            result = _extract_from_blocks(content)
            if result:
                return result

        # Output with candidates (generateContent embedded)
        output = step.get("output", {})
        for candidate in output.get("candidates", []):
            parts = candidate.get("content", {}).get("parts", [])
            result = _extract_from_parts(parts)
            if result:
                return result

    # Format 3: Top-level candidates (old generateContent fallback)
    for candidate in response_data.get("candidates", []):
        parts = candidate.get("content", {}).get("parts", [])
        result = _extract_from_parts(parts)
        if result:
            return result

    return None


def _validate_generated_image(path: str, min_size_bytes: int = 1024) -> bool:
    """Check if a generated image is valid (not empty, not corrupted)."""
    try:
        if not os.path.isfile(path):
            return False
        fsize = os.path.getsize(path)
        if fsize < min_size_bytes:
            logger.warning("Generated image too small (%d bytes): %s", fsize, path)
            return False
        # Try opening to verify it's a valid image
        from PIL import Image as PILImage
        img = PILImage.open(path)
        img.verify()
        return True
    except Exception:
        logger.warning("Generated image validation failed: %s", path, exc_info=True)
        return False


async def _retry_image_generation(
    gen_func,
    prompt: str,
    seg_id: int,
    note_id: int,
    max_retries: int = 2,
    **kwargs,
) -> str | None:
    """Retry image generation up to max_retries times on failure."""
    for attempt in range(max_retries + 1):
        try:
            result = await gen_func(prompt=prompt, seg_id=seg_id, note_id=note_id, **kwargs)
            if result and _validate_generated_image(result):
                if attempt > 0:
                    logger.info("Image generation succeeded on retry %d for seg %s", attempt, seg_id)
                return result
            if attempt < max_retries:
                logger.warning("Retry %d/%d for seg %s (invalid or empty image)", attempt + 1, max_retries, seg_id)
                await asyncio.sleep(1.5)
        except Exception:
            if attempt < max_retries:
                logger.warning("Retry %d/%d for seg %s after exception", attempt + 1, max_retries, seg_id, exc_info=True)
                await asyncio.sleep(1.5)
    return None


async def _generate_image_with_free_model_wrapper(
    prompt: str,
    seg_id: int,
    note_id: int,
    **kwargs,
) -> str | None:
    """Wrapper to match _retry_image_generation signature."""
    return await _generate_image_with_free_model(
        prompt=prompt, seg_id=seg_id, note_id=note_id, **kwargs,
    )


async def _generate_image_with_free_model(
    prompt: str,
    seg_id: int,
    note_id: int,
    model_config: dict | None = None,
    width: int = 1920,
    height: int = 1080,
    prefix: str = "bg",
) -> str | None:
    """Try to generate image using free models, based on model_config.model_id."""
    cfg = model_config or {}
    model_id = cfg.get("model_id", "")

    if not model_id or model_id not in FREE_IMAGE_MODELS:
        return None

    model_info = FREE_IMAGE_MODELS[model_id]
    provider = model_info["provider"]
    model_name = model_info["model_name"]

    # Clamp dimensions to model's max
    w = min(width, model_info["max_width"])
    h = min(height, model_info["max_height"])

    if provider == "pollinations":
        return await _generate_image_pollinations(
            prompt=prompt, width=w, height=h,
            seg_id=seg_id, note_id=note_id,
            model_name=model_name, prefix=prefix,
        )
    elif provider == "gemini":
        return await _generate_image_gemini(
            prompt=prompt, width=w, height=h,
            seg_id=seg_id, note_id=note_id, prefix=prefix,
            model_config=cfg,
        )

    return None


# ---------------------------------------------------------------------------
# LLM-powered prompt generation
# ---------------------------------------------------------------------------

async def _generate_slide_prompts_llm(
    segments: list[dict],
    style: str,
    image_style: str | None = None,
    llm_config: dict | None = None,
) -> dict[int, str]:
    """Use LLM to generate high-quality image prompts for each slide background."""
    cfg = llm_config or {}
    api_key = cfg.get("api_key") or settings.LLM_API_KEY
    base_url = cfg.get("base_url") or settings.LLM_BASE_URL
    model_id = cfg.get("model_id") or settings.LLM_MODEL

    if not api_key:
        logger.info("No LLM API key for slide prompt generation, using fallback")
        return {}

    seg_list = ""
    for seg in segments:
        seg_id = seg.get("id", 0)
        text = seg.get("text", "")
        visual_hint = seg.get("visual_hint", "")
        emotion = seg.get("emotion", "neutral")
        seg_list += f"\n[{seg_id}] 旁白：{text} | 画面提示：{visual_hint} | 情绪：{emotion}"

    style_desc = {
        "knowledge": "教育科普风格，专业大气",
        "story": "电影叙事风格，有故事感",
        "checklist": "信息图风格，简洁有力",
    }.get(style, "高质量演示文稿风格")

    if image_style and image_style in SLIDE_STYLE_MAP:
        style_desc += f"，画面风格：{SLIDE_STYLE_MAP[image_style]}"

    system_msg = f"""You are an expert AI image prompt engineer specializing in presentation and social media backgrounds.

Your task: for each script segment below, write a high-quality English prompt for AI image generation models (Flux/Stable Diffusion).

CRITICAL RULES:
- ALL prompts MUST be in English (Flux works best with English)
- Each prompt must be 50-120 words long with rich visual detail
- Describe specific subjects, composition, lighting, color palette, mood, texture
- The background must leave visual space for overlaid text (not too busy in center)
- Include EXACTLY: {NEGATIVE_PROMPT}
- Include quality markers: {QUALITY_SUFFIX}
- Aspect ratio context: 16:9 horizontal for PPT, 3:4 vertical for social media
- Avoid abstract concepts — describe concrete visual scenes
- Use professional photography/film terminology for realistic styles

Style direction: {style_desc}

Output PURE JSON only, no markdown, no explanation:
{{"1": "highly detailed prompt in English...", "2": "another prompt...", ...}}

Key is segment ID, value is the complete English prompt."""


    user_msg = f"请为以下{len(segments)}段旁白生成对应的PPT背景图提示词：{seg_list}"

    try:
        http_client = httpx.AsyncClient(verify=False) if "coze" in str(base_url) else None
        client = AsyncOpenAI(api_key=api_key, base_url=base_url, http_client=http_client)
        try:
            response = await client.chat.completions.create(
                model=model_id,
                max_tokens=2048,
                temperature=0.7,
                messages=[
                    {"role": "system", "content": system_msg},
                    {"role": "user", "content": user_msg},
                ],
            )
            text = response.choices[0].message.content or ""
        finally:
            if http_client:
                await http_client.aclose()
    except Exception:
        logger.exception("LLM call for slide prompts failed")
        return {}

    parsed = None
    try:
        parsed = json.loads(text)
    except json.JSONDecodeError:
        start = text.find("{")
        end = text.rfind("}")
        if start != -1 and end != -1:
            try:
                parsed = json.loads(text[start:end + 1])
            except json.JSONDecodeError:
                pass

    if not parsed or not isinstance(parsed, dict):
        logger.warning("Failed to parse LLM slide prompt response")
        return {}

    result = {}
    for seg in segments:
        seg_id = seg.get("id", 0)
        key = str(seg_id)
        if key in parsed and isinstance(parsed[key], str):
            result[seg_id] = parsed[key]
    return result


# ---------------------------------------------------------------------------
# Prompt fallback
# ---------------------------------------------------------------------------

def _build_slide_prompt(segment: dict, style: str, image_style: str | None = None) -> str:
    """Build a high-quality slide background prompt from segment data (fallback)."""
    visual_hint = segment.get("visual_hint", "")
    emotion = segment.get("emotion", "neutral")
    style_prefix = {
        "knowledge": "professional presentation background, clean modern design, educational atmosphere",
        "story": "cinematic storytelling scene, narrative composition, emotional depth",
        "checklist": "minimalist infographic background, structured layout, clean aesthetic",
    }.get(style, "high quality professional presentation background")
    emotion_suffix = {
        "inspiring": "uplifting golden hour light, hopeful atmosphere, warm sun rays",
        "positive": "bright warm natural light, cheerful inviting mood, sunny atmosphere",
        "neutral": "calm balanced atmosphere, soft diffused ambient lighting, professional",
        "calm": "peaceful serene atmosphere, gentle diffused morning light, tranquil",
        "melancholy": "somber reflective mood, cool muted blue tones, soft rain aesthetic",
        "dramatic": "intense dramatic atmosphere, strong contrast chiaroscuro lighting",
        "energetic": "vibrant dynamic energy, bold vivid saturated colors, motion blur",
        "reflective": "contemplative introspective mood, soft ethereal backlight, thoughtful",
    }.get(emotion, "balanced atmosphere, professional studio lighting")

    parts = [style_prefix]
    if image_style and image_style in SLIDE_STYLE_MAP:
        parts.append(SLIDE_STYLE_MAP[image_style])
    if visual_hint:
        parts.append(visual_hint)
    parts.append(emotion_suffix)
    parts.append(QUALITY_SUFFIX)
    parts.append(NEGATIVE_PROMPT)
    parts.append("16:9 horizontal composition, suitable for presentation background with text overlay space")
    return ", ".join(parts)


# ---------------------------------------------------------------------------
# Fallback gradient background image generation
# ---------------------------------------------------------------------------

def _create_gradient_image(palette: dict, w: int = 1920, h: int = 1080) -> str:
    """Create a gradient PNG and return its path."""
    img = Image.new("RGB", (w, h))
    draw = ImageDraw.Draw(img)
    top = palette["top"]
    bottom = palette["bottom"]
    for y in range(h):
        ratio = y / h
        r = int(top[0] + (bottom[0] - top[0]) * ratio)
        g = int(top[1] + (bottom[1] - top[1]) * ratio)
        b = int(top[2] + (bottom[2] - top[2]) * ratio)
        draw.line([(0, y), (w, y)], fill=(r, g, b))

    tmp = str(SLIDES_DIR / f"_gradient_{id(palette)}.png")
    img.save(tmp, "PNG")
    return tmp


# ---------------------------------------------------------------------------
# Darken / preprocess AI background for text readability
# ---------------------------------------------------------------------------

def _preprocess_bg_image(path: str) -> str:
    """Subtly enhance AI background for text readability without destroying detail."""
    try:
        img = Image.open(path).convert("RGBA")
        img = img.resize((1920, 1080), Image.LANCZOS)
        # Lighter dark overlay — preserve more of the original image
        dark = Image.new("RGBA", (1920, 1080), (0, 0, 0, 50))
        img = Image.alpha_composite(img, dark)
        # Very subtle blur — sharp enough to keep AI detail visible
        img = img.filter(ImageFilter.GaussianBlur(radius=0.6))
        out = path.replace(".png", "_processed.png")
        img.convert("RGB").save(out, "PNG", quality=95)
        return out
    except Exception:
        logger.exception("Failed to preprocess bg image %s", path)
        return path


# ---------------------------------------------------------------------------
# Text splitting
# ---------------------------------------------------------------------------

def _split_heading_body(text: str, emotion: str) -> tuple[str, str]:
    """Split segment text into a heading and body for a slide."""
    for sep in ["。", "：", "：", "—", "，"]:
        idx = text.find(sep)
        if 0 < idx < 40:
            heading = text[:idx + 1].strip()
            body = text[idx + 1:].strip()
            if body:
                return heading, body

    if len(text) <= 30:
        return text, ""

    if emotion:
        labels = {
            "inspiring": "启 发", "positive": "积 极", "neutral": "要 点",
            "calm": "平 静", "melancholy": "感 悟", "dramatic": "重 点",
            "energetic": "能 量", "reflective": "思 考",
        }
        return labels.get(emotion, emotion), text

    return "要 点", text


# ---------------------------------------------------------------------------
# PPTX building
# ---------------------------------------------------------------------------

def _set_shape_transparency(shape, alpha_pct: int):
    """Set transparency on a filled shape. alpha_pct: 0-100 (0=opaque, 100=invisible)."""
    from lxml import etree
    from pptx.oxml.ns import qn
    spPr = shape._element.find(qn('p:spPr'))
    if spPr is None:
        return
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is None:
        return
    srgbClr = solidFill.find(qn('a:srgbClr'))
    if srgbClr is None:
        return
    # Remove existing alpha
    for existing in srgbClr.findall(qn('a:alpha')):
        srgbClr.remove(existing)
    alpha_elem = etree.SubElement(srgbClr, qn('a:alpha'))
    alpha_elem.set('val', str(int((100 - alpha_pct) * 1000)))  # val in 1/1000ths of percent


def _build_pptx(
    title: str,
    subtitle: str,
    content_slides: list[dict],
    output_path: str,
) -> str:
    """Build a .pptx file with title + content slides.

    content_slides: list of dicts with keys:
        heading, body, bg_path (str|None), page_num, total_pages
    """
    prs = Presentation()
    prs.slide_width = PPT_WIDTH_EMU
    prs.slide_height = PPT_HEIGHT_EMU

    # ---- Blank layout ----
    blank_layout = prs.slide_layouts[6]  # usually "Blank"

    # ============ Title Slide ============
    slide = prs.slides.add_slide(blank_layout)

    # Background image or solid fill
    title_data = content_slides[0] if content_slides else {}
    title_bg = title_data.get("bg_path")

    if title_bg and os.path.isfile(title_bg):
        slide.shapes.add_picture(
            title_bg,
            Emu(0), Emu(0),
            prs.slide_width, prs.slide_height,
        )
        # Semi-transparent dark overlay rectangle
        overlay = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Emu(0), Emu(0),
            prs.slide_width, prs.slide_height,
        )
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0x10, 0x20, 0x40)
        overlay.line.fill.background()
        _set_shape_transparency(overlay, 60)  # 60% transparent
    else:
        # Dark blue background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x0F, 0x20, 0x4B)

    # Accent bar at top
    accent_bar = slide.shapes.add_shape(1, Emu(0), Emu(0), prs.slide_width, Emu(55000))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = RGBColor(0x3B, 0x82, 0xF6)
    accent_bar.line.fill.background()

    # Title text
    title_left = IN(1.2)
    title_top = IN(2.2)
    title_width = IN(10.9)
    title_height = IN(2.0)
    txBox = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name = PPTX_FONT_NAME

    # Subtitle
    if subtitle:
        sub_left = IN(1.2)
        sub_top = IN(4.4)
        sub_width = IN(10.9)
        sub_height = IN(1.0)
        txBox2 = slide.shapes.add_textbox(sub_left, sub_top, sub_width, sub_height)
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(22)
        run2.font.color.rgb = RGBColor(0xDC, 0xDC, 0xDC)
        run2.font.name = PPTX_FONT_NAME

    # Decorative line under title
    line_left = IN(4.5)
    line_top = IN(4.2)
    line_width = IN(4.3)
    line_height = Emu(22000)
    line_shape = slide.shapes.add_shape(1, line_left, line_top, line_width, line_height)
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = RGBColor(0x3B, 0x82, 0xF6)
    line_shape.line.fill.background()

    # ============ Content Slides ============
    for cs in content_slides:
        slide = prs.slides.add_slide(blank_layout)

        # Background
        bg_path = cs.get("bg_path")
        if bg_path and os.path.isfile(bg_path):
            slide.shapes.add_picture(
                bg_path,
                Emu(0), Emu(0),
                prs.slide_width, prs.slide_height,
            )
            # Dark overlay
            overlay = slide.shapes.add_shape(
                1, Emu(0), Emu(0),
                prs.slide_width, prs.slide_height,
            )
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(0x08, 0x10, 0x20)
            overlay.line.fill.background()
            _set_shape_transparency(overlay, 65)  # 65% transparent
        else:
            # Gradient-style solid dark background
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0x14, 0x28, 0x50)

        # Accent bar at top
        accent = slide.shapes.add_shape(1, Emu(0), Emu(0), prs.slide_width, Emu(55000))
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(0x3B, 0x82, 0xF6)
        accent.line.fill.background()

        # Heading
        heading = cs.get("heading", "")
        body = cs.get("body", "")
        page_num = cs.get("page_num", 0)
        total_pages = cs.get("total_pages", 0)

        if heading:
            h_left = IN(0.8)
            h_top = IN(0.6)
            h_width = IN(11.5)
            h_height = IN(1.2)
            txBox = slide.shapes.add_textbox(h_left, h_top, h_width, h_height)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = heading
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.font.name = PPTX_FONT_NAME

            # Separator line under heading
            sep = slide.shapes.add_shape(1, h_left, IN(1.8), IN(1.8), Emu(28000))
            sep.fill.solid()
            sep.fill.fore_color.rgb = RGBColor(0x3B, 0x82, 0xF6)
            sep.line.fill.background()

        # Body text
        if body:
            b_left = IN(0.8)
            b_top = IN(2.2)
            b_width = IN(11.5)
            b_height = IN(4.5)
            txBox = slide.shapes.add_textbox(b_left, b_top, b_width, b_height)
            tf = txBox.text_frame
            tf.word_wrap = True

            # Wrap long body text
            max_chars = 38
            wrapped = textwrap.fill(body, width=max_chars)
            lines = wrapped.split("\n")

            for i, line in enumerate(lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.space_after = Pt(6)
                run = p.add_run()
                run.text = line
                run.font.size = Pt(22)
                run.font.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
                run.font.name = PPTX_FONT_NAME

        # Page number
        if total_pages > 0:
            pn_left = IN(11.5)
            pn_top = IN(6.8)
            pn_width = IN(1.5)
            pn_height = IN(0.4)
            txBox = slide.shapes.add_textbox(pn_left, pn_top, pn_width, pn_height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            run = p.add_run()
            run.text = f"{page_num} / {total_pages}"
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(0xA0, 0xA0, 0xA0)
            run.font.name = PPTX_FONT_NAME

    prs.save(output_path)
    return output_path


# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------

async def generate_slides(
    script: dict,
    note_id: int,
    image_style: str | None = None,
    model_config: dict | None = None,
    llm_config: dict | None = None,
) -> dict:
    """Generate a PPT file with AI backgrounds from a script.

    Returns a dict with keys: local_path, page_count, has_ai_bg.
    """
    _ensure_dir()

    segments = script.get("segments", [])
    title = script.get("title", "")
    style = script.get("style", "knowledge")

    # --- Step 1: LLM prompts ---
    llm_prompts = await _generate_slide_prompts_llm(segments, style, image_style, llm_config)

    # --- Step 2: Generate AI background images (parallel) ---
    bg_tasks = []
    seg_ids = []
    for seg in segments:
        seg_id = seg.get("id", 0)
        seg_ids.append(seg_id)
        prompt = llm_prompts.get(seg_id) or _build_slide_prompt(seg, style, image_style)
        bg_tasks.append(_generate_slide_bg_image(prompt, seg_id, note_id, model_config))

    # Title background
    title_prompt = (
        f"professional presentation cover background, abstract elegant design, "
        f"{SLIDE_STYLE_MAP.get(image_style or 'realistic', '')}, "
        f"no text, dramatic lighting, 16:9 landscape"
    )
    bg_tasks.append(_generate_slide_bg_image(title_prompt, -1, note_id, model_config))

    bg_results = await asyncio.gather(*bg_tasks, return_exceptions=True)

    # Log any errors (403 should have been caught by pre-check)
    for i, result in enumerate(bg_results):
        if isinstance(result, Exception):
            logger.warning("Background generation task %d failed: %s", i, result)

    # Map seg_id → bg_path
    bg_paths: dict[int, str | None] = {}
    for seg_id, result in zip(seg_ids, bg_results[:len(seg_ids)]):
        bg_paths[seg_id] = result if isinstance(result, str) else None

    title_bg_path = bg_results[-1] if isinstance(bg_results[-1], str) else None

    # --- Step 3: Preprocess backgrounds (darken for text readability) ---
    processed_bg: dict[int, str | None] = {}
    for seg_id, raw_path in bg_paths.items():
        if raw_path and os.path.isfile(raw_path):
            processed_bg[seg_id] = _preprocess_bg_image(raw_path)
        else:
            # Create gradient fallback
            pal = GRADIENT_PALETTES[(seg_id + 1) % len(GRADIENT_PALETTES)]
            processed_bg[seg_id] = _create_gradient_image(pal)

    processed_title_bg = None
    if title_bg_path and os.path.isfile(title_bg_path):
        processed_title_bg = _preprocess_bg_image(title_bg_path)
    else:
        processed_title_bg = _create_gradient_image(GRADIENT_PALETTES[0])

    # --- Step 4: Build content_slides data ---
    style_label = {
        "knowledge": "知识解读",
        "story": "故事讲述",
        "checklist": "清单体",
    }.get(style, "")

    content_slides = []

    # Title slide (first entry, special handling)
    content_slides.append({
        "heading": "",
        "body": "",
        "bg_path": processed_title_bg,
        "page_num": 0,
        "total_pages": 0,
    })

    # Content slides
    for seg in segments:
        seg_id = seg.get("id", 0)
        text = seg.get("text", "").strip()
        if not text:
            continue
        heading, body = _split_heading_body(text, seg.get("emotion", ""))
        content_slides.append({
            "heading": heading,
            "body": body,
            "bg_path": processed_bg.get(seg_id),
            "page_num": len(content_slides),
            "total_pages": len(segments),
        })

    # Update page numbers (title is page 0, content starts at 1)
    for i, cs in enumerate(content_slides):
        cs["page_num"] = i
        cs["total_pages"] = len(content_slides) - 1  # exclude title

    # --- Step 5: Build .pptx ---
    output_path = str(SLIDES_DIR / f"note_{note_id}.pptx")
    subtitle = f"{len(segments)} 页 · {style_label}"

    _build_pptx(
        title=title,
        subtitle=subtitle,
        content_slides=content_slides,
        output_path=output_path,
    )

    has_ai_bg = any(bg_paths.values())

    return {
        "local_path": output_path,
        "page_count": len(content_slides),
        "has_ai_bg": has_ai_bg,
    }


# ===========================================================================
# 小红书图文生成 (3:4 vertical, 1080x1440, swipe carousel)
# ===========================================================================

XHS_W = 1080
XHS_H = 1440

# 小红书 warm / aesthetic color palettes
XHS_PALETTES = [
    {"top": (255, 240, 245), "bottom": (255, 220, 230)},   # 蜜桃粉
    {"top": (255, 248, 235), "bottom": (255, 230, 190)},   # 奶茶橙
    {"top": (235, 245, 255), "bottom": (200, 225, 255)},   # 天空蓝
    {"top": (240, 255, 240), "bottom": (200, 240, 210)},   # 抹茶绿
    {"top": (245, 240, 255), "bottom": (225, 210, 255)},   # 香芋紫
    {"top": (255, 255, 240), "bottom": (255, 245, 200)},   # 奶油黄
]

XHS_ACCENT = (233, 69, 96)  # 小红书红


# --- Font metrics: CJK glyph width ≈ font_size px ---

def _measure_text_width(text: str, font) -> int:
    """Accurately measure pixel width of text with given font."""
    dummy = Image.new("RGBA", (1, 1))
    draw = ImageDraw.Draw(dummy)
    bbox = draw.textbbox((0, 0), text, font=font)
    return bbox[2] - bbox[0]


def _wrap_text_to_width(text: str, font, max_width: int) -> list[str]:
    """Wrap text to fit within max_width pixels, breaking at CJK chars or spaces."""
    if not text:
        return []
    lines: list[str] = []
    current = ""
    for ch in text:
        test = current + ch
        if _measure_text_width(test, font) > max_width and current:
            lines.append(current)
            current = ch
        else:
            current = test
    if current:
        lines.append(current)
    return lines


def _split_xhs_text(text: str) -> tuple[str, str]:
    """Split segment text for 小红书 layout.

    Strategy: take the first natural phrase as heading (short, punchy),
    rest as body. Unlike PPT we want headings that are SHORT —
    ideally 6-15 chars for visual impact.
    """
    # Try to split at first punctuation, but cap heading to ~20 chars
    for sep in ["。", "，", "！", "？", "：", "—", "；"]:
        idx = text.find(sep)
        if 0 < idx <= 18:
            heading = text[:idx + 1].strip()
            body = text[idx + 1:].strip()
            if body:
                return heading, body

    # If short enough, use as heading only
    if len(text) <= 18:
        return text, ""

    # Force break at ~18 chars
    return text[:18], text[18:]


def _render_xhs_cover(
    title: str,
    subtitle: str,
    bg_image_path: str | None = None,
    palette: dict | None = None,
    out_path: str | None = None,
) -> dict | None:
    """Render a 小红书 cover page — bold title, rich visual layering, mobile-optimized."""
    try:
        is_ai_bg = bg_image_path and os.path.isfile(bg_image_path)

        # ==================================================================
        # 1. BACKGROUND — AI image with warm overlay, or gradient with texture
        # ==================================================================
        if is_ai_bg:
            img = Image.open(bg_image_path).convert("RGBA")
            img = img.resize((XHS_W, XHS_H), Image.LANCZOS)
            # Warm dark gradient overlay (darker at bottom for text readability)
            overlay = Image.new("RGBA", (XHS_W, XHS_H), (0, 0, 0, 0))
            odraw = ImageDraw.Draw(overlay)
            for y in range(XHS_H):
                ratio = y / XHS_H
                # Center stays brighter, edges darken
                alpha = int(60 + 40 * ratio)
                r, g, b = int(15 + 5 * ratio), int(8 + 2 * ratio), int(8 + 2 * ratio)
                odraw.line([(0, y), (XHS_W, y)], fill=(r, g, b, alpha))
            img = Image.alpha_composite(img, overlay)
            # Edge vignette
            vignette = Image.new("RGBA", (XHS_W, XHS_H), (0, 0, 0, 0))
            vd = ImageDraw.Draw(vignette)
            for i in range(80):
                a = int(4 * i / 80)
                vd.rounded_rectangle(
                    [(i, i), (XHS_W - i, XHS_H - i)],
                    radius=max(0, 40 - i // 2),
                    outline=(0, 0, 0, a), width=1,
                )
            img = Image.alpha_composite(img, vignette)
        else:
            pal = palette or XHS_PALETTES[0]
            img = Image.new("RGBA", (XHS_W, XHS_H))
            draw = ImageDraw.Draw(img)
            top_c, bot_c = pal["top"], pal["bottom"]
            # Smooth gradient
            for y in range(XHS_H):
                r = y / XHS_H
                c = tuple(int(top_c[i] + (bot_c[i] - top_c[i]) * r) for i in range(3))
                draw.line([(0, y), (XHS_W, y)], fill=(*c, 255))
            # Texture: scattered translucent circles for depth
            import random
            random.seed(42)
            for _ in range(80):
                dx = random.randint(40, XHS_W - 40)
                dy = random.randint(40, XHS_H - 40)
                r = random.randint(12, 36)
                a = random.randint(8, 22)
                draw.ellipse([(dx, dy), (dx + r, dy + r)], fill=(*top_c, a))
            # Large subtle geometric shapes
            random.seed(7)
            for _ in range(6):
                dx = random.randint(60, XHS_W - 260)
                dy = random.randint(60, XHS_H - 260)
                s = random.randint(140, 280)
                a = random.randint(6, 16)
                draw.rounded_rectangle(
                    [(dx, dy), (dx + s, dy + s)],
                    radius=60, fill=(*top_c, a),
                )

        draw = ImageDraw.Draw(img)

        # ==================================================================
        # 2. DECORATIVE ELEMENTS — layered geometric accents
        # ==================================================================
        # Top accent gradient bar
        for y in range(6):
            alpha = 180 - y * 25
            draw.line([(0, y), (XHS_W, y)], fill=(*XHS_ACCENT, max(20, alpha)))

        # Top-left: overlapping semi-transparent circles
        for offset, alpha, r in [
            (-60, 35, 200), (-30, 25, 140), (60, 18, 100),
        ]:
            draw.ellipse(
                [(offset, -80 + offset), (offset + r, -80 + offset + r)],
                fill=(*XHS_ACCENT, alpha),
            )
        # Top-right: subtle accent dot
        draw.ellipse(
            [(XHS_W - 130, 60), (XHS_W - 40, 150)],
            fill=(*XHS_ACCENT, 30),
        )

        # Bottom decorative elements
        # Thick-to-thin horizontal accent lines
        draw.rounded_rectangle(
            [(70, XHS_H - 150), (XHS_W - 70, XHS_H - 146)],
            radius=2, fill=(255, 255, 255, 35),
        )
        draw.rounded_rectangle(
            [(140, XHS_H - 140), (XHS_W - 140, XHS_H - 137)],
            radius=2, fill=(255, 255, 255, 20),
        )

        # ==================================================================
        # 3. TITLE — large, centered, with shadow and accent underline
        # ==================================================================
        title_font = _load_font(72)
        max_title_w = XHS_W - 180
        title_lines = _wrap_text_to_width(title, title_font, max_title_w)

        line_h = 90
        total_title_h = len(title_lines) * line_h
        y_start = max(160, (XHS_H - total_title_h) // 2 - 80)

        for i, line in enumerate(title_lines):
            lw = _measure_text_width(line, title_font)
            x = (XHS_W - lw) // 2
            y = y_start + i * line_h
            # Text shadow (deeper for better readability)
            draw.text((x + 4, y + 4), line, fill=(0, 0, 0, 80), font=title_font)
            draw.text((x + 2, y + 2), line, fill=(0, 0, 0, 40), font=title_font)
            # Main text
            text_color = (255, 255, 255, 255) if is_ai_bg else (40, 30, 30, 255)
            draw.text((x, y), line, fill=text_color, font=title_font)

        # Accent underline — two-tone, wider, centered
        line_y = y_start + len(title_lines) * line_h + 20
        last_line_w = _measure_text_width(title_lines[-1] if title_lines else title, title_font)
        line_w = min(last_line_w, 320)
        line_x = (XHS_W - line_w) // 2
        # Left half (accent red)
        half = line_w // 2
        draw.rounded_rectangle(
            [(line_x, line_y), (line_x + half, line_y + 5)],
            radius=3, fill=(*XHS_ACCENT, 240),
        )
        # Right half (lighter)
        draw.rounded_rectangle(
            [(line_x + half, line_y), (line_x + line_w, line_y + 5)],
            radius=3, fill=(*XHS_ACCENT, 130),
        )

        # ==================================================================
        # 4. SUBTITLE — larger, with decorative diamonds on sides
        # ==================================================================
        if subtitle:
            sub_font = _load_font(32)
            sw = _measure_text_width(subtitle, sub_font)
            sx = (XHS_W - sw) // 2
            sy = line_y + 26
            sub_color = (240, 235, 230, 210) if is_ai_bg else (120, 110, 105, 255)
            sub_shadow = (0, 0, 0, 40) if is_ai_bg else (0, 0, 0, 0)
            # Decorative diamond shapes
            diamond_gap = 40
            for side_x in [sx - diamond_gap, sx + sw + diamond_gap]:
                diamond_cy = sy + 12
                pts = [
                    (side_x, diamond_cy - 6),   # top
                    (side_x + 6, diamond_cy),    # right
                    (side_x, diamond_cy + 6),   # bottom
                    (side_x - 6, diamond_cy),    # left
                ]
                draw.polygon(pts, fill=(*sub_color[:3], 150))
            # Subtitle text
            if sub_shadow[3]:
                draw.text((sx + 2, sy + 2), subtitle, fill=sub_shadow, font=sub_font)
            draw.text((sx, sy), subtitle, fill=sub_color, font=sub_font)

        # ==================================================================
        # 5. SWIPE HINT — larger, with animated-style arrows
        # ==================================================================
        hint_font = _load_font(26)
        hint_text = "← 左右滑动查看更多 →"
        hw = _measure_text_width(hint_text, hint_font)
        hx = (XHS_W - hw) // 2
        hy = XHS_H - 100
        hint_color = (210, 205, 200, 180) if is_ai_bg else (160, 155, 150, 200)
        # Small double-arrow decorations flanking the text
        arr_y = hy + 10
        for side, dx in [("left", -1), ("right", 1)]:
            arr_x = (hx - 30) if side == "left" else (hx + hw + 10)
            draw.polygon(
                [(arr_x + dx * 16, arr_y), (arr_x + dx * 4, arr_y - 8), (arr_x + dx * 4, arr_y + 8)],
                fill=(*hint_color[:3], 100),
            )
            draw.polygon(
                [(arr_x + dx * 20, arr_y), (arr_x + dx * 8, arr_y - 8), (arr_x + dx * 8, arr_y + 8)],
                fill=(*hint_color[:3], 60),
            )
        draw.text((hx, hy), hint_text, fill=hint_color, font=hint_font)

        # ==================================================================
        # Finalize
        # ==================================================================
        img_rgb = Image.new("RGB", (XHS_W, XHS_H), (255, 255, 255))
        img_rgb.paste(img, mask=img.split()[3])
        if out_path:
            img_rgb.save(out_path, "PNG", quality=95)
            return {"local_path": out_path}
        return {"image": img_rgb}
    except Exception:
        logger.exception("Failed to render XHS cover")
        return None


def _render_xhs_content(
    heading: str,
    body: str,
    page_num: int,
    total_pages: int,
    bg_image_path: str | None = None,
    palette: dict | None = None,
    out_path: str | None = None,
) -> dict | None:
    """Render a 小红书 content page — spacious card layout, warm tones, mobile-optimized."""
    try:
        is_ai_bg = bg_image_path and os.path.isfile(bg_image_path)

        # ==================================================================
        # 1. BACKGROUND
        # ==================================================================
        if is_ai_bg:
            img = Image.open(bg_image_path).convert("RGBA")
            img = img.resize((XHS_W, XHS_H), Image.LANCZOS)
            # Softer, warmer dark overlay
            overlay = Image.new("RGBA", (XHS_W, XHS_H), (0, 0, 0, 0))
            odraw = ImageDraw.Draw(overlay)
            for y in range(XHS_H):
                ratio = y / XHS_H
                alpha = int(50 + 15 * ratio)
                odraw.line([(0, y), (XHS_W, y)], fill=(12, 7, 7, alpha))
            img = Image.alpha_composite(img, overlay)
            img = img.filter(ImageFilter.GaussianBlur(radius=1.0))
        else:
            pal = palette or XHS_PALETTES[0]
            img = Image.new("RGBA", (XHS_W, XHS_H))
            draw = ImageDraw.Draw(img)
            top_c, bot_c = pal["top"], pal["bottom"]
            for y in range(XHS_H):
                r = y / XHS_H
                c = tuple(int(top_c[i] + (bot_c[i] - top_c[i]) * r) for i in range(3))
                draw.line([(0, y), (XHS_W, y)], fill=(*c, 255))
            # Soft texture
            import random
            random.seed(page_num * 13)
            for _ in range(40):
                dx = random.randint(30, XHS_W - 30)
                dy = random.randint(30, XHS_H - 30)
                s = random.randint(16, 40)
                a = random.randint(6, 18)
                draw.rounded_rectangle(
                    [(dx, dy), (dx + s, dy + s)],
                    radius=s // 2, fill=(*top_c, a),
                )

        # ==================================================================
        # 2. TOP ACCENT BAR — gradient from red to transparent
        # ==================================================================
        bar_draw = ImageDraw.Draw(img)
        for y in range(5):
            alpha = 200 - y * 35
            bar_draw.line([(0, y), (XHS_W, y)], fill=(*XHS_ACCENT, max(30, alpha)))
        # Thin secondary bar below
        bar_draw.rectangle(
            [(0, 6), (XHS_W, 8)],
            fill=(*XHS_ACCENT, 40),
        )

        # ==================================================================
        # 3. CORNER DECORATIONS — subtle elegance
        # ==================================================================
        # Top-left corner arcs
        d = ImageDraw.Draw(img)
        for offset, alpha in [(20, 25), (35, 18), (50, 10)]:
            d.arc(
                [(offset, offset), (offset + 80, offset + 80)],
                start=180, end=270, fill=(*XHS_ACCENT, alpha), width=2,
            )
        # Top-right corner arcs
        for offset, alpha in [(20, 25), (35, 18), (50, 10)]:
            d.arc(
                [(XHS_W - offset - 80, offset), (XHS_W - offset, offset + 80)],
                start=270, end=360, fill=(*XHS_ACCENT, alpha), width=2,
            )

        # ==================================================================
        # 4. MAIN CONTENT CARD — larger, positioned higher
        # ==================================================================
        card_margin = 50
        card_top = 260
        card_bottom = XHS_H - 100

        card = Image.new("RGBA", (XHS_W, XHS_H), (0, 0, 0, 0))
        card_draw = ImageDraw.Draw(card)
        # Softer shadow
        card_draw.rounded_rectangle(
            [(card_margin + 4, card_top + 5), (XHS_W - card_margin + 4, card_bottom + 5)],
            radius=32, fill=(0, 0, 0, 22),
        )
        # Shadow feather
        card_draw.rounded_rectangle(
            [(card_margin + 8, card_top + 10), (XHS_W - card_margin + 8, card_bottom + 10)],
            radius=32, fill=(0, 0, 0, 12),
        )
        # Main card body
        card_draw.rounded_rectangle(
            [(card_margin, card_top), (XHS_W - card_margin, card_bottom)],
            radius=32, fill=(255, 255, 255, 242),
        )
        # Subtle card border
        card_draw.rounded_rectangle(
            [(card_margin, card_top), (XHS_W - card_margin, card_bottom)],
            radius=32, outline=(255, 255, 255, 80), width=1,
        )
        img = Image.alpha_composite(img, card)
        draw = ImageDraw.Draw(img)

        # Card-top accent stripe (gradient two-tone)
        stripe_left = card_margin + 28
        stripe_right = XHS_W - card_margin - 28
        stripe_w = stripe_right - stripe_left
        mid = stripe_left + stripe_w // 2
        stripe_y = card_top + 18
        draw.rounded_rectangle(
            [(stripe_left, stripe_y), (mid, stripe_y + 6)],
            radius=3, fill=(*XHS_ACCENT, 210),
        )
        draw.rounded_rectangle(
            [(mid, stripe_y), (stripe_right, stripe_y + 6)],
            radius=3, fill=(*XHS_ACCENT, 100),
        )

        # ==================================================================
        # 5. NUMBER BADGE — floating circle on top of card
        # ==================================================================
        badge_r = 34
        badge_cx = XHS_W // 2
        badge_cy = card_top
        # Shadow ring
        draw.ellipse(
            [(badge_cx - badge_r + 3, badge_cy - badge_r + 4),
             (badge_cx + badge_r + 3, badge_cy + badge_r + 4)],
            fill=(0, 0, 0, 25),
        )
        # Badge fill
        draw.ellipse(
            [(badge_cx - badge_r, badge_cy - badge_r),
             (badge_cx + badge_r, badge_cy + badge_r)],
            fill=(*XHS_ACCENT, 255),
        )
        # Badge outer ring
        draw.ellipse(
            [(badge_cx - badge_r, badge_cy - badge_r),
             (badge_cx + badge_r, badge_cy + badge_r)],
            outline=(255, 255, 255, 120), width=2,
        )
        # Number
        num_font = _load_font(28)
        num_text = str(page_num)
        nw = _measure_text_width(num_text, num_font)
        draw.text(
            (badge_cx - nw // 2, badge_cy - 15),
            num_text, fill=(255, 255, 255, 255), font=num_font,
        )

        # ==================================================================
        # 6. TEXT AREA — larger fonts, warmer color, generous spacing
        # ==================================================================
        text_left = card_margin + 48
        text_right = XHS_W - card_margin - 48
        text_width = text_right - text_left
        y_cursor = card_top + 52

        # Heading
        if heading:
            heading_font = _load_font(46)
            heading_lines = _wrap_text_to_width(heading, heading_font, text_width)
            heading_color = (45, 35, 30, 255)  # warm dark brown
            heading_line_h = 62
            for hl in heading_lines:
                if y_cursor + heading_line_h > card_bottom - 30:
                    break
                draw.text((text_left, y_cursor), hl, fill=heading_color, font=heading_font)
                y_cursor += heading_line_h

            # Accent separator — wider, two-tone
            sep_w = 64
            draw.rounded_rectangle(
                [(text_left, y_cursor + 8), (text_left + sep_w, y_cursor + 13)],
                radius=3, fill=(*XHS_ACCENT, 220),
            )
            draw.rounded_rectangle(
                [(text_left + sep_w + 12, y_cursor + 9), (text_left + sep_w + 40, y_cursor + 12)],
                radius=2, fill=(*XHS_ACCENT, 80),
            )
            y_cursor += 34

        # Body text
        if body:
            body_font = _load_font(38)
            body_line_h = 54
            body_lines = _wrap_text_to_width(body, body_font, text_width)
            body_color = (75, 65, 60, 255)  # warm dark gray
            for bi, bl in enumerate(body_lines):
                if y_cursor + body_line_h > card_bottom - 36:
                    # Overflow indicator
                    dot_font = _load_font(38)
                    draw.text(
                        (text_left, y_cursor), "···",
                        fill=(180, 170, 165, 255), font=dot_font,
                    )
                    break
                # First line emphasis (slightly bold-looking via color)
                line_color = (55, 40, 35, 255) if bi == 0 and len(body_lines) > 2 else body_color
                draw.text((text_left, y_cursor), bl, fill=line_color, font=body_font)
                y_cursor += body_line_h

        # ==================================================================
        # 7. PAGE INDICATOR DOTS — more prominent
        # ==================================================================
        dot_y = XHS_H - 56
        dot_spacing = 22
        total_dots_w = total_pages * dot_spacing
        dot_start_x = (XHS_W - total_dots_w) // 2
        for i in range(total_pages):
            cx = dot_start_x + i * dot_spacing + dot_spacing // 2
            is_current = i + 1 == page_num
            if is_current:
                # Active dot — filled + glow ring
                draw.ellipse(
                    [(cx - 8, dot_y - 8), (cx + 8, dot_y + 8)],
                    fill=(*XHS_ACCENT, 255),
                )
                draw.ellipse(
                    [(cx - 11, dot_y - 11), (cx + 11, dot_y + 11)],
                    outline=(*XHS_ACCENT, 100), width=2,
                )
            else:
                draw.ellipse(
                    [(cx - 5, dot_y - 5), (cx + 5, dot_y + 5)],
                    fill=(200, 190, 185, 170),
                )

        # ==================================================================
        # Finalize
        # ==================================================================
        img_rgb = Image.new("RGB", (XHS_W, XHS_H), (255, 255, 255))
        img_rgb.paste(img, mask=img.split()[3])
        if out_path:
            img_rgb.save(out_path, "PNG", quality=95)
            return {"local_path": out_path}
        return {"image": img_rgb}
    except Exception:
        logger.exception("Failed to render XHS content page")
        return None


async def generate_xhs_images(
    script: dict,
    note_id: int,
    image_style: str | None = None,
    model_config: dict | None = None,
    llm_config: dict | None = None,
) -> dict:
    """Generate 小红书 swipe images (3:4 vertical, 1080x1440)."""
    import zipfile

    _ensure_dir()

    segments = script.get("segments", [])
    title = script.get("title", "")
    style = script.get("style", "knowledge")

    # --- LLM prompts ---
    llm_prompts = await _generate_slide_prompts_llm(segments, style, image_style, llm_config)

    # --- Resolve model_id for prompt adaptation ---
    _cfg = model_config or {}
    _model_id = _cfg.get("model_id") or settings.IMAGE_MODEL

    # --- AI backgrounds ---
    bg_tasks = []
    seg_ids = []
    for seg in segments:
        seg_id = seg.get("id", 0)
        seg_ids.append(seg_id)
        prompt = llm_prompts.get(seg_id) or _build_xhs_prompt(seg, style, image_style, _model_id)
        bg_tasks.append(_generate_xhs_bg_image(prompt, seg_id, note_id, model_config))

    cover_prompt = (
        f"xiaohongshu style cover background, aesthetic warm tones, "
        f"{SLIDE_STYLE_MAP.get(image_style or 'realistic', '')}, "
        f"no text, dreamy soft lighting, 3:4 portrait, lifestyle photography"
    )
    bg_tasks.append(_generate_xhs_bg_image(cover_prompt, -1, note_id, model_config))

    bg_results = await asyncio.gather(*bg_tasks, return_exceptions=True)

    # Log any errors (403 should have been caught by pre-check, but handle gracefully)
    for i, result in enumerate(bg_results):
        if isinstance(result, Exception):
            logger.warning("Background generation task %d failed: %s", i, result)

    bg_paths: dict[int, str | None] = {}
    for seg_id, result in zip(seg_ids, bg_results[:len(seg_ids)]):
        bg_paths[seg_id] = result if isinstance(result, str) else None
    cover_bg_path = bg_results[-1] if isinstance(bg_results[-1], str) else None

    # --- Render ---
    style_label = {
        "knowledge": "知识解读", "story": "故事讲述", "checklist": "清单体",
    }.get(style, "")

    image_paths = []
    page_index = 0

    # Cover
    cover_path = str(SLIDES_DIR / f"note_{note_id}_xhs_{page_index:03d}.png")
    cover_result = _render_xhs_cover(
        title=title,
        subtitle=f"{len(segments)} 图 · {style_label}",
        bg_image_path=cover_bg_path,
        palette=XHS_PALETTES[0],
        out_path=cover_path,
    )
    if cover_result and cover_result.get("local_path"):
        image_paths.append(cover_result["local_path"])
        page_index += 1

    # Content pages
    for seg in segments:
        seg_id = seg.get("id", 0)
        text = seg.get("text", "").strip()
        if not text:
            continue

        heading, body = _split_xhs_text(text)
        bg_path = bg_paths.get(seg_id)
        palette = XHS_PALETTES[(seg_id) % len(XHS_PALETTES)]

        img_path = str(SLIDES_DIR / f"note_{note_id}_xhs_{page_index:03d}.png")
        result = _render_xhs_content(
            heading=heading,
            body=body,
            page_num=page_index,
            total_pages=len(segments),
            bg_image_path=bg_path,
            palette=palette,
            out_path=img_path,
        )
        if result and result.get("local_path"):
            image_paths.append(result["local_path"])
            page_index += 1

    # --- Pack into zip ---
    zip_path = str(SLIDES_DIR / f"note_{note_id}_xhs.zip")
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for p in image_paths:
            if os.path.isfile(p):
                zf.write(p, os.path.basename(p))

    has_ai_bg = any(bg_paths.values())

    return {
        "local_path": zip_path,
        "page_count": len(image_paths),
        "has_ai_bg": has_ai_bg,
        "image_paths": image_paths,
    }


async def _generate_xhs_bg_image(
    prompt: str,
    seg_id: int,
    note_id: int,
    model_config: dict | None = None,
) -> str | None:
    """Generate a 3:4 portrait background image for 小红书.

    Priority: free models (Pollinations / Gemini) -> SiliconFlow paid models -> None.
    """
    cfg = model_config or {}
    model_id = cfg.get("model_id") or settings.IMAGE_MODEL

    # --- Try free models first (with retry) ---
    if _is_free_model(model_id):
        result = await _retry_image_generation(
            _generate_image_with_free_model_wrapper,
            prompt=prompt, seg_id=seg_id, note_id=note_id,
            max_retries=2, model_config=model_config,
            width=768, height=1024, prefix="xhs_bg",
        )
        if result:
            return result
        logger.warning("Free model %s failed for XHS segment %s, using gradient fallback", model_id, seg_id)
        return None

    # --- Existing SiliconFlow / paid model path ---
    api_key = cfg.get("api_key") or settings.IMAGE_API_KEY
    base_url = cfg.get("base_url") or settings.IMAGE_BASE_URL

    if not api_key:
        return None

    # Determine optimal inference steps per model
    m = model_id.lower()
    if "schnell" in m:
        steps = 4
    elif "lightning" in m:
        steps = 4
    elif "flux" in m and "dev" in m:
        steps = 28
    elif "flux" in m:
        steps = 20
    elif "kolors" in m:
        steps = 25
    elif "qwen-image" in m or "z-image" in m or "ernie" in m:
        steps = 20
    elif "sd3" in m or "stable-diffusion-3" in m:
        steps = 28
    elif "turbo" in m:
        steps = 8
    elif "playground" in m:
        steps = 25
    else:
        steps = 20

    try:
        url = f"{base_url.rstrip('/')}/images/generations"
        headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
        payload = {
            "model": model_id,
            "prompt": prompt,
            "image_size": "768x1024",
            "num_inference_steps": steps,
        }

        async with httpx.AsyncClient(timeout=120.0) as client:
            resp = await client.post(url, json=payload, headers=headers)
            if resp.status_code == 403:
                logger.error("Image model %s returned 403 Forbidden — account balance insufficient or model requires payment. Model ID: %s, Base URL: %s", model_id, model_id, base_url)
                raise HTTPException(status_code=403, detail=f"图片模型 {model_id} 返回 403：余额不足或该模型需要付费，请充值后重试或切换到免费模型（Kolors）")
            resp.raise_for_status()
            data = resp.json()

        image_url = None
        images = data.get("images", [])
        if images and isinstance(images, list):
            image_url = images[0].get("url")
        elif data.get("data"):
            image_url = data["data"][0].get("url")

        if not image_url:
            return None

        out_path = SLIDES_DIR / f"note_{note_id}_xhs_bg_{seg_id}.png"
        async with httpx.AsyncClient(timeout=60.0) as client:
            img_resp = await client.get(image_url)
            img_resp.raise_for_status()
            out_path.write_bytes(img_resp.content)

        return str(out_path)
    except HTTPException:
        raise  # Let 403 errors bubble up to caller
    except Exception:
        logger.exception("Failed to generate XHS background for segment %s", seg_id)
        return None


def _build_xhs_prompt(segment: dict, style: str, image_style: str | None = None, model_id: str = "") -> str:
    """Build a 小红书 style background prompt, adapted to the image model."""
    visual_hint = segment.get("visual_hint", "")
    emotion = segment.get("emotion", "neutral")
    style_prefix = {
        "knowledge": "xiaohongshu aesthetic lifestyle photo, clean informative composition, warm inviting tones",
        "story": "xiaohongshu cinematic storytelling, emotional portrait mood, dreamy atmospheric lighting",
        "checklist": "xiaohongshu organized infographic background, modern pastel aesthetic, neat layout",
    }.get(style, "xiaohongshu aesthetic photography, warm and inviting lifestyle mood")
    emotion_suffix = {
        "inspiring": "golden hour sunlight, hopeful warm glow, inspirational atmosphere",
        "positive": "bright cheerful daylight, sun-kissed warm tones, happy vibes",
        "neutral": "soft diffused natural lighting, calm balanced composition",
        "calm": "gentle morning light through window, peaceful serenity, tranquil mood",
        "melancholy": "moody soft overcast lighting, cool muted blue-grey palette",
        "dramatic": "dramatic side lighting, rich contrast, intense atmosphere",
        "energetic": "vibrant bold saturated colors, dynamic energetic composition",
        "reflective": "soft ethereal backlight, contemplative introspective mood",
    }.get(emotion, "soft natural window lighting, lifestyle photography aesthetic")

    parts = [style_prefix]
    if image_style and image_style in SLIDE_STYLE_MAP:
        parts.append(SLIDE_STYLE_MAP[image_style])
    if visual_hint:
        parts.append(visual_hint)
    parts.append(emotion_suffix)
    parts.append(QUALITY_SUFFIX)
    parts.append(NEGATIVE_PROMPT)

    # Model-specific prompt enhancements
    m = model_id.lower()
    if "kolors" in m:
        parts.append("小红书风格, 生活方式摄影, 暖色调, 精致构图, 高画质")
        parts.append("no text, no letters, no watermark, 3:4 portrait")
    elif "qwen-image" in m or "z-image" in m or "ernie" in m:
        parts.append("小红书风格, 生活方式摄影, 暖色调, 精致构图, 高画质, 无文字, 无水印")
    elif "flux" in m:
        parts.append("3:4 vertical portrait, xiaohongshu social media aesthetic, lifestyle photography, 4k detailed")
    elif "turbo" in m:
        parts.append("3:4 portrait, social media aesthetic, lifestyle photo, detailed")
    else:
        parts.append("3:4 portrait, xiaohongshu aesthetic, lifestyle photography, detailed")

    return ", ".join(parts)


# ============================================================
# Template-based XHS prompt generation (no LLM needed)
# ============================================================

# Topic → visual scene mapping for XHS content
_XHS_TOPIC_SCENES: dict[str, str] = {
    # Lifestyle & daily
    "穿搭": "fashion outfit flat lay on wooden table, natural sunlight, aesthetic clothing arrangement, magazine style",
    "美妆": "beauty cosmetics on marble vanity, soft glam lighting, elegant skincare bottles arrangement",
    "护肤": "skincare routine on bathroom shelf, morning light, fresh clean aesthetic, spa atmosphere",
    "美食": "food photography overhead, warm restaurant lighting, artful plating, steam rising, cozy dining atmosphere",
    "咖啡": "coffee cup on rustic wooden table, morning sunlight through window, latte art, cozy cafe vibes",
    "旅行": "travel destination landscape, golden hour lighting, wanderlust aesthetic, scenic viewpoint",
    "家居": "cozy home interior, natural light through curtains, minimalist decor, hygge atmosphere, house plants",
    "健身": "fitness lifestyle, morning workout at golden hour, clean gym aesthetic, motivational atmosphere",
    "读书": "book and coffee on cozy reading nook, soft warm lamp light, autumn vibes, quiet afternoon",
    "摄影": "camera on desk, creative workspace, warm ambient light, artistic photography gear arrangement",
    "音乐": "headphones on desk, moody lighting, vinyl record aesthetic, creative music atmosphere",
    "手帐": "bullet journal flat lay, washi tape and stickers, warm desk lamp, creative stationery aesthetic",
    "插花": "fresh flowers arrangement, natural window light, botanical aesthetic, soft pastel tones",
    "宠物": "cute pet in cozy home, soft natural lighting, heartwarming atmosphere, lifestyle pet photography",
    "办公": "minimal desk setup, warm ambient lighting, productive workspace aesthetic, modern home office",
    # Knowledge & education
    "职场": "professional workspace with warm lighting, modern office aesthetic, career inspiration",
    "效率": "minimal productivity desk setup, clean organized workspace, morning focus atmosphere",
    "学习": "study desk with warm lamp light, books and notes, cozy academic atmosphere, focused ambiance",
    "心理": "calm meditation space, soft diffused lighting, peaceful serenity, mindfulness atmosphere",
    "理财": "minimal financial planning desk, warm golden light, organized aesthetic, prosperity vibes",
    "成长": "sunrise over mountain path, hopeful golden light, personal growth journey, inspirational vista",
    # Technology
    "科技": "sleek tech product on minimal desk, blue ambient glow, futuristic clean design, modern aesthetic",
    "数码": "gadget unboxing flat lay, clean white background, modern tech lifestyle, soft studio lighting",
    # Generic fallbacks
    "生活": "lifestyle flat lay on warm wooden surface, morning sunlight, cozy aesthetic, soft shadows",
    "分享": "aesthetic mood board, warm tones, creative inspiration, soft natural lighting",
    "推荐": "product display on aesthetic background, soft studio lighting, modern lifestyle, clean composition",
    "日常": "daily life vignette, golden hour light through window, warm cozy atmosphere, candid moments",
    "测评": "product review setup, clean studio lighting, detailed close-up, professional product photography",
    "经验": "notebook and coffee, soft morning light, contemplative atmosphere, knowledge sharing vibes",
    "教程": "step by step creative process, clean organized workspace, bright natural light, tutorial atmosphere",
}

# Emotion → lighting/mood mapping
_XHS_EMOTION_ATMOSPHERE: dict[str, str] = {
    "inspiring": "golden hour sunlight, hopeful warm glow, inspirational atmosphere, uplifting",
    "positive": "bright cheerful daylight, sun-kissed warm tones, happy energetic vibes",
    "neutral": "soft diffused natural lighting, calm balanced composition, serene",
    "calm": "gentle morning light through window, peaceful serenity, tranquil mood, quiet",
    "melancholy": "moody soft overcast lighting, cool muted tones, quiet contemplative atmosphere",
    "dramatic": "dramatic side lighting, rich contrast, intense cinematic atmosphere",
    "energetic": "vibrant bold saturated colors, dynamic energetic composition, lively",
    "reflective": "soft ethereal backlight, contemplative introspective mood, dreamy",
    "warm": "warm golden sunset light, cozy intimate atmosphere, soft bokeh background",
    "fresh": "fresh morning dew, crisp clean lighting, bright airy atmosphere, spring vibes",
}

# Style keywords for different image styles (XHS-optimized)
_XHS_STYLE_KEYWORDS: dict[str, str] = {
    "realistic": "photorealistic, natural colors, soft natural lighting, shallow depth of field, 85mm lens look",
    "illustration": "flat vector illustration, modern clean style, soft pastel colors, editorial illustration",
    "anime": "anime art style, cel shading, warm color palette, Japanese aesthetic, Studio Ghibli inspired",
    "watercolor": "watercolor painting style, soft washes, gentle color blending, artistic paper texture",
    "comic": "webtoon style, clean line art, soft cell shading, modern manhwa aesthetic",
    "minimal": "minimalist design, clean lines, ample negative space, muted pastel palette, modern Nordic aesthetic",
    "cinematic": "cinematic composition, dramatic lighting, rich color grading, depth of field, film grain",
    "dreamy": "soft focus, dreamy ethereal glow, pastel color palette, bokeh highlights, romantic atmosphere",
    "vintage": "film photography aesthetic, warm retro tones, light leaks, nostalgic 90s vibes, grain texture",
    "cyberpunk": "neon lights urban night, cyberpunk aesthetic, purple-blue lighting, futuristic Tokyo vibes",
    "nature": "lush botanical garden, natural sunlight through leaves, organic textures, biophilic design",
    "abstract": "abstract fluid art, soft color gradients, organic flowing shapes, modern art gallery aesthetic",
}


def _detect_topic(text: str) -> str:
    """Detect the main topic from text content by keyword matching."""
    for keyword, scene in _XHS_TOPIC_SCENES.items():
        if keyword in text:
            return scene
    return _XHS_TOPIC_SCENES["生活"]  # default


def _infer_emotion_atmosphere(text: str) -> str:
    """Infer emotion/atmosphere from text sentiment cues."""
    positive_words = ["喜欢", "爱", "开心", "快乐", "幸福", "美好", "推荐", "值得", "很棒", "超好", "安利", "种草"]
    calm_words = ["安静", "平静", "放松", "治愈", "温柔", "慢慢", "宁静", "静静", "平和"]
    energetic_words = ["激情", "燃", "炸裂", "爆", "冲", "燃爆", "热血", "惊艳", "绝了"]
    fresh_words = ["清新", "春天", "新鲜", "清晨", "活力", "元气", "朝气"]

    positive_count = sum(1 for w in positive_words if w in text)
    calm_count = sum(1 for w in calm_words if w in text)
    energetic_count = sum(1 for w in energetic_words if w in text)
    fresh_count = sum(1 for w in fresh_words if w in text)

    if energetic_count > 0:
        return _XHS_EMOTION_ATMOSPHERE["energetic"]
    if fresh_count > positive_count and fresh_count > 0:
        return _XHS_EMOTION_ATMOSPHERE["fresh"]
    if calm_count > positive_count:
        return _XHS_EMOTION_ATMOSPHERE["calm"]
    if positive_count > 0:
        return _XHS_EMOTION_ATMOSPHERE["positive"]
    return _XHS_EMOTION_ATMOSPHERE["warm"]  # default warm for XHS


def _build_xhs_prompts_from_template(
    title: str,
    segments: list[dict],
    image_style: str | None = None,
    model_id: str = "",
) -> list[str]:
    """Build XHS image prompts from content using templates — no LLM needed.

    Returns a list of prompts, one per segment (plus cover).
    """
    prompts = []

    # Resolve style keywords
    style_keywords = ""
    if image_style and image_style in _XHS_STYLE_KEYWORDS:
        style_keywords = _XHS_STYLE_KEYWORDS[image_style]
    elif image_style and image_style in SLIDE_STYLE_MAP:
        style_keywords = SLIDE_STYLE_MAP[image_style]

    # Model-specific postfix
    m = model_id.lower()
    if "flux" in m:
        model_postfix = "3:4 vertical portrait, xiaohongshu social media aesthetic, lifestyle photography, 4k detailed"
    elif "turbo" in m:
        model_postfix = "3:4 portrait, social media aesthetic, lifestyle photo, detailed"
    elif "gemini" in m:
        model_postfix = "xiaohongshu vertical post 3:4 ratio, lifestyle photography, soft natural lighting"
    else:
        model_postfix = "3:4 portrait, xiaohongshu aesthetic, lifestyle photography, detailed"

    # --- Cover prompt: based on title ---
    topic_scene = _detect_topic(title)
    emotion = _infer_emotion_atmosphere(title)
    cover_prompt_parts = [
        "xiaohongshu cover image aesthetic",
        topic_scene,
        "warm inviting atmosphere",
        emotion,
        "elegant composition with copy space for text overlay in center",
        "soft blurred background edges for text readability",
    ]
    if style_keywords:
        cover_prompt_parts.insert(1, style_keywords)
    cover_prompt_parts.append(QUALITY_SUFFIX)
    cover_prompt_parts.append(NEGATIVE_PROMPT)
    cover_prompt_parts.append(model_postfix)
    prompts.append(", ".join(cover_prompt_parts))

    # --- Content page prompts ---
    for seg in segments:
        text = seg.get("text", "")
        visual_hint = seg.get("visual_hint", "")
        emotion_key = seg.get("emotion", "warm")

        topic_scene = _detect_topic(text)
        atmosphere = _XHS_EMOTION_ATMOSPHERE.get(emotion_key, _XHS_EMOTION_ATMOSPHERE["warm"])

        page_prompt_parts = [
            "xiaohongshu content image background",
            topic_scene,
            atmosphere,
            "soft focus edges for text overlay readability",
        ]
        if style_keywords:
            page_prompt_parts.insert(1, style_keywords)
        if visual_hint:
            page_prompt_parts.insert(2, visual_hint)

        page_prompt_parts.append(QUALITY_SUFFIX)
        page_prompt_parts.append(NEGATIVE_PROMPT)
        page_prompt_parts.append(model_postfix)
        prompts.append(", ".join(page_prompt_parts))

    return prompts
